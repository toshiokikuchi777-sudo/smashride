from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define colors
    NEON_BLUE = RGBColor(0, 229, 255)
    DARK_BG = RGBColor(30, 30, 30)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(50, 50, 50)
    TABLE_HEADER_BG = RGBColor(40, 40, 40)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_img_placeholder(slide, left, top, width, height):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = GRAY
        shape.line.color.rgb = WHITE
        txBox = slide.shapes.add_textbox(left, top + height/2 - Inches(0.1), width, Inches(0.2))
        p = txBox.text_frame.paragraphs[0]
        p.text = "[ IMAGE ]"
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = WHITE
        p.font.size = Pt(8)

    def create_table_slide(prs, title_text):
        slide_layout = prs.slide_layouts[5] # Blank
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, DARK_BG)
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].font.color.rgb = NEON_BLUE
        title_shape.text_frame.paragraphs[0].font.size = Pt(18)
        title_shape.text_frame.paragraphs[0].font.bold = True
        return slide

    def add_styled_table(slide, rows, cols, left, top, width, height, data):
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Set column widths
        first_col_width = width * 0.3
        table.columns[0].width = int(first_col_width)
        for i in range(1, cols):
            table.columns[i].width = int((width - first_col_width) / (cols - 1))

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                # Header row styling
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_HEADER_BG
                
                # Set text and style
                p = cell.text_frame.paragraphs[0]
                p.text = data[r][c]
                p.font.size = Pt(8.5 if r > 0 else 9.5)
                p.font.color.rgb = NEON_BLUE if c == 0 or r == 0 else WHITE
                p.font.bold = (r == 0 or c == 0)
                cell.margin_left = cell.margin_right = Pt(5)
                cell.margin_top = cell.margin_bottom = Pt(2)
        
        return table_shape

    # 1. 概要 & インスピレーション
    s1 = create_table_slide(prs, "01. 概要と企画の背景")
    data1 = [
        ["項目", "内容 / 詳細"],
        ["タイトル", "SMASH RIDE 🛹"],
        ["コンセプト", "スケボーのスピード感 × 缶破壊の爽快感"],
        ["主なターゲット", "Robloxユーザー (8〜18歳)、カジュアル層"],
        ["Drill Digging", "成長モデル。訪問数: 9,700万+ | リリース: 2024/02\nLink: https://www.roblox.com/games/15873911365/"],
        ["Cut Trees", "対象物の変化。訪問数: 3,400万+ | リリース: 2025/07\nLink: https://www.roblox.com/games/135880624242201/"],
        ["Pet Simulator 99", "収集・倍率UP。訪問数: 23億+ | リリース: 2023/12\nLink: https://www.roblox.com/games/5969560376/"],
        ["クリック/中毒性", "シンプルな操作、数値インフレ、報酬サイクルの高速化"]
    ]
    # Adjust table height to fit more text
    add_styled_table(s1, len(data1), 2, Inches(0.3), Inches(0.9), Inches(5.8), Inches(5.5), data1)
    add_img_placeholder(s1, Inches(6.3), Inches(1.0), Inches(3.4), Inches(5.5))

    # 2. コアループ & ステージ
    s2 = create_table_slide(prs, "02. コアループと進行")
    data2 = [
        ["フェーズ", "アクション ＆ 報酬詳細"],
        ["行動 (Action)", "スケボーでの高速移動 ＆ ハンマーによる缶の破砕"],
        ["報酬 (Reward)", "SCRAPポイント獲得、統計(破壊数)の加算"],
        ["強化 (Next)", "ペットガチャ(倍率UP)、ハンマー購入(範囲・性能)"],
        ["Lv1 - 5", "初期ワールド(赤/青)、基本的な操作と循環を習得"],
        ["Lv6 - 10", "成長期(紫/黄)、マスター装備を目指してSCRAPを蓄積"],
        ["目標 (Goal)", "次のワールドへのロック解除、ランキング上位挑戦"]
    ]
    add_styled_table(s2, len(data2), 2, Inches(0.3), Inches(1.0), Inches(6.0), Inches(4.0), data2)
    add_img_placeholder(s2, Inches(6.5), Inches(1.0), Inches(3.2), Inches(5.5))

    # 3. 装備システム (スケボー & ハンマー)
    s3 = create_table_slide(prs, "03. 装備システム：ルールと性能設計")
    data3 = [
        ["要素 / 装備", "特徴 ＆ アップグレード効果"],
        ["スケボー", "アイテム化、Lv1(16速)→Lv5(32速)への高速化"],
        ["ジャンプ力", "Lv1(50)→Lv5(100)、範囲破壊に直接影響"],
        ["Smash Action", "ジャンプ着地で範囲破壊、ボーナス1.5倍、スキル性"],
        ["BASIC / SHOCK", "初心者用(赤のみ) / 衝撃波(青まで破壊可能)"],
        ["MULTI / HYBRID", "複数色倍率(緑) / 万能範囲(紫まで対応)"],
        ["MASTER", "究極装備。全色対応、圧倒的な破壊速度とスケール"]
    ]
    add_styled_table(s3, len(data3), 2, Inches(0.3), Inches(1.0), Inches(6.5), Inches(4.5), data3)
    add_img_placeholder(s3, Inches(7.0), Inches(1.0), Inches(2.7), Inches(5.5))

    # 4. ワールド ＆ ガチャ
    s4 = create_table_slide(prs, "04. ワールド展開 ＆ ペットシステム")
    data4 = [
        ["ワールド名", "テーマ ＆ 解放条件"],
        ["現代都市", "スタート地点。缶・ビル景観。初期解放"],
        ["SF未来", "サイバーパンク、ホバー、ネオン缶。ステージ10"],
        ["古代遺跡", "壺・石像破壊。累計5,000缶破壊で解放"],
        ["水中/マグマ", "貝殻/溶岩。過酷環境、累計1〜2万破壊"],
        ["ペット要素", "11種、3スロット。ポイント倍率の飛躍的向上"],
        ["ガチャの進化", "物理的なガチャゾーンへ移行(体験の向上)"]
    ]
    add_styled_table(s4, len(data4), 2, Inches(0.3), Inches(1.0), Inches(6.0), Inches(4.5), data4)
    add_img_placeholder(s4, Inches(6.5), Inches(1.0), Inches(3.2), Inches(5.5))

    # 5. マネタイズ ＆ システム設計
    s5 = create_table_slide(prs, "05. システム ＆ ビジネス設計")
    data5 = [
        ["項目", "内容 ＆ 期待される効果"],
        ["勝利条件", "終了なきエンドレス成長、ランキング上位到達"],
        ["失敗条件", "ミスやGame Over排除。純粋な癒やしと成長を提供"],
        ["成長パス", "ハンマー範囲・倍率、ペット倍率、スケボー機動力"],
        ["GamePass", "装備枠拡張、永久倍率ブースト、高ランク装備解放"],
        ["DevProducts", "SCRAPパック、ガチャ連、限定ペット販売"],
        ["広告報酬", "無料ガチャ権、一時的なスピード/倍率、ブースト時間"]
    ]
    add_styled_table(s5, len(data5), 2, Inches(0.3), Inches(1.0), Inches(6.5), Inches(4.5), data5)
    add_img_placeholder(s5, Inches(7.0), Inches(1.0), Inches(2.7), Inches(5.5))

    # 6. 進捗 ＆ ロードマップ
    s6 = create_table_slide(prs, "06. 開発進捗 ＆ 未来の展望")
    data6 = [
        ["フェーズ", "マイルストーン ＆ 進捗状況"],
        ["実装済み", "物理破壊、11種ペット、ガチャ、DataStore同期"],
        ["現状課題", "MASTER後のインフレ、Region3負荷対策、PvP要否"],
        ["Phase 1", "チュートリアル、パス設定、バランス微調整、リリース"],
        ["Phase 2", "SF未来ワールド、ボスバトル、アチーブメント"],
        ["KPI目標", "DAU 1k、プレイ時間30min、Day1継続 40%"],
        ["ビジョン", "Robloxにおける新感覚アクションクリッカーの確立"]
    ]
    add_styled_table(s6, len(data6), 2, Inches(0.3), Inches(1.0), Inches(5.5), Inches(4.5), data6)
    add_img_placeholder(s6, Inches(6.0), Inches(1.0), Inches(3.7), Inches(5.5))

    prs.save('Smash_Ride_Table_Deck.pptx')
    print("Table Presentation created successfully: Smash_Ride_Table_Deck.pptx")

if __name__ == "__main__":
    create_presentation()
